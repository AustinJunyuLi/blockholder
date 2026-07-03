"""Business-format slidepack: Liquidity, Activism Disclosure, and Takeover Premia.

Generates pres/blockholder_seminar_40min.pptx — a 40-minute business-styled
deck (decision-oriented headlines, figure-led, technical material in backups)
sharing its narrative spine with the academic Beamer deck
(quality_reports/plans/2026-06-11_slidepack-40min-design.md).

Design system: strategy-consulting chrome — navy full-bleed cover and divider,
agenda, tracked small-caps kickers, white cards with colored top edges, stat
callouts, hairline tables. Display mathematics is typeset with the
manuscript's Computer Modern (eq_render.py: xelatex + preview -> transparent
600-dpi PNGs) and placed as pictures, so equations match the figures.

Figures are rasterized from the canonical PDFs (pres/figures + empirics/output)
to 300-dpi PNGs in pres/pptx_assets/ via pdftocairo. Fact 1/2 numbers are read
from empirics/output/*.csv at build time, so re-running after the event-study
refresh updates the evidence slides automatically.

Usage:
    .venv/bin/python pres/make_pptx.py
"""

from __future__ import annotations

import csv
import os
import subprocess

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

import eq_render

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
ASSETS = os.path.join(HERE, "pptx_assets")
OUT_PPTX = os.path.join(HERE, "blockholder_seminar_40min.pptx")

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
EMU_IN = 914400

# ── palette: strategy-consulting navy/charcoal chrome; Paul Tol muted
#    accents retained where they key to the figure layer ────────────────────
INK = RGBColor(0x21, 0x21, 0x21)      # charcoal body text
GREY = RGBColor(0x5A, 0x5A, 0x5A)     # support text
FAINT = RGBColor(0xF3, 0xF4, 0xF6)    # soft band fill (cool light grey)
HAIR = RGBColor(0xD5, 0xD9, 0xDE)     # hairline rules / card borders
NAVY = RGBColor(0x1F, 0x3A, 0x5F)     # primary accent (kickers, bands)
NAVY_D = RGBColor(0x16, 0x29, 0x3F)   # deep navy (cover base)
ICE = RGBColor(0xC9, 0xD6, 0xE8)      # light text on navy
BLUE = NAVY                           # structural accent alias (legacy name)
DATABLUE = RGBColor(0x44, 0x77, 0xAA)  # series blue used inside figures
ROSE = RGBColor(0xEE, 0x66, 0x77)
TEAL = RGBColor(0x44, 0xAA, 0x99)
SAND = RGBColor(0xDD, 0xCC, 0x77)
CYAN = RGBColor(0x88, 0xCC, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Helvetica Neue"
FOOT_TITLE = "LIQUIDITY · ACTIVISM DISCLOSURE · TAKEOVER PREMIA"

MARGIN = Inches(0.60)
CONTENT_W = Inches(12.13)

# figure PDFs to rasterize: (source path relative to ROOT, asset name)
# Shared manuscript figures come from numerical_output/ (canonical); the four
# slide-only variants from pres/figures/ (pyfig.slide_figures).
FIGS = [
    ("numerical_output/fig_nonmonotone.pdf", "nonmonotone"),
    ("numerical_output/fig_decomposition.pdf", "decomposition"),
    ("numerical_output/fig_wedge_primitives.pdf", "wedge"),
    ("numerical_output/fig_ge_decomposition.pdf", "ge_decomp"),
    ("pres/figures/fig_disclosure_slopes.pdf", "slopes"),
    ("numerical_output/fig_welfare.pdf", "welfare"),
    ("numerical_output/fig_cutoff_structure.pdf", "cutoffs"),
    ("numerical_output/fig_cutoffs_kappa.pdf", "cutoffs_kappa"),
    ("pres/figures/fig_sensitivity_panel1.pdf", "sens1"),
    ("pres/figures/fig_sensitivity_panel2.pdf", "sens2"),
    ("pres/figures/fig_noisy_rumor.pdf", "rumor"),
    ("empirics/output/fact1_delay.pdf", "fact1"),
    ("empirics/output/fact2_car.pdf", "fact2"),
]


def build_assets() -> None:
    os.makedirs(ASSETS, exist_ok=True)
    for rel, name in FIGS:
        src = os.path.join(ROOT, rel)
        dst = os.path.join(ASSETS, name)
        if not os.path.exists(src):
            print(f"  !! missing {rel} (slide will show a placeholder)")
            continue
        subprocess.run(
            ["pdftocairo", "-png", "-r", "300", "-singlefile", src, dst],
            check=True,
        )
    print(f"assets -> {ASSETS}")


# ── data hooks: evidence numbers read from the empirics outputs ─────────────

def fact1_numbers() -> dict:
    path = os.path.join(ROOT, "empirics/output/fact1_summary.csv")
    out = {"pre_med": "7", "post_med": "5", "pre_5bd": "36%", "post_5bd": "76%"}
    try:
        with open(path) as fh:
            for row in csv.DictReader(fh):
                if row["window"].startswith("pre"):
                    out["pre_med"] = f"{float(row['median']):.0f}"
                    out["pre_5bd"] = f"{100 * float(row['share_within_5bd']):.0f}%"
                else:
                    out["post_med"] = f"{float(row['median']):.0f}"
                    out["post_5bd"] = f"{100 * float(row['share_within_5bd']):.0f}%"
    except OSError:
        pass
    return out


def fact2_numbers() -> dict:
    """Headline coefficients from spec 2 on CAR[-1,+1]; em-dash until run."""
    path = os.path.join(ROOT, "empirics/output/fact2_regressions.csv")
    out = {"beta": "—", "beta_t": "—", "delta": "—", "delta_t": "—",
           "amihud": "—", "amihud_t": "—", "n": "—"}
    try:
        with open(path) as fh:
            for row in csv.DictReader(fh):
                if (row["spec"] == "2_interaction"
                        and row["depvar"] == "car_m1_p1"):
                    if row["param"] == "post":
                        out["beta"] = f"{float(row['coef']) * 100:+.2f}pp"
                        out["beta_t"] = f"t={float(row['t']):.2f}"
                        out["n"] = row["n"]
                    if row["param"] == "post_x_lnamihud":
                        out["delta"] = f"{float(row['coef']) * 100:+.2f}pp"
                        out["delta_t"] = f"t={float(row['t']):.2f}"
                    if row["param"] == "ln_amihud":
                        out["amihud"] = f"{float(row['coef']) * 100:+.2f}pp"
                        out["amihud_t"] = f"t={float(row['t']):.2f}"
    except OSError:
        pass
    return out


# ── slide toolkit ───────────────────────────────────────────────────────────

def _track(run, spc: int = 160) -> None:
    """Letter-spacing in 1/100 pt (OOXML rPr@spc) — tracked small caps."""
    run._r.get_or_add_rPr().set("spc", str(spc))


def _flat(shp) -> None:
    """Kill all themed effects (drop shadows) on a shape or connector.

    shadow.inherit = False writes an empty <a:effectLst/>, but renderers
    other than PowerPoint (LibreOffice, Google Slides) still apply the
    theme effect referenced by the shape's <p:style>, so remove it too.
    """
    el = shp._element
    style = el.find(qn("p:style"))
    if style is not None:
        el.remove(style)
    shp.shadow.inherit = False


def _segments(txt: str, opt: dict):
    """'**bold**' segments inside a string; returns (text, bold, color)."""
    base_bold = opt.get("bold", False)
    parts = txt.split("**")
    out = []
    for i, seg in enumerate(parts):
        if seg == "":
            continue
        out.append((seg, base_bold or (i % 2 == 1), None))
    return out or [("", base_bold, None)]


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.blank = self.prs.slide_layouts[6]
        self.n = 0
        self.eqs = eq_render.ensure_equations(ASSETS)

    # ── chrome ──────────────────────────────────────────────────────────────

    def slide(self, kicker: str = "", headline: str = "", footer: bool = True):
        s = self.prs.slides.add_slide(self.blank)
        self.n += 1
        if kicker:
            sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(0.455),
                                    Inches(0.10), Inches(0.10))
            sq.fill.solid(); sq.fill.fore_color.rgb = NAVY
            sq.line.fill.background(); _flat(sq)
            tb = self._text(s, Inches(0.80), Inches(0.345), Inches(11), Inches(0.3))
            p = tb.paragraphs[0]
            r = p.add_run(); r.text = kicker.upper()
            f = r.font; f.name = FONT; f.size = Pt(10.5); f.bold = True
            f.color.rgb = NAVY
            _track(r, 170)
        if headline:
            tb = self._text(s, Inches(0.58), Inches(0.66), CONTENT_W, Inches(1.0))
            p = tb.paragraphs[0]
            r = p.add_run(); r.text = headline
            f = r.font; f.name = FONT; f.size = Pt(26); f.bold = True
            f.color.rgb = INK
        if footer:
            self.hairline(s, MARGIN, Inches(7.06), CONTENT_W,
                          color=RGBColor(0xE3, 0xE6, 0xEA))
            tb = self._text(s, MARGIN, Inches(7.13), Inches(9), Inches(0.28))
            p = tb.paragraphs[0]
            r = p.add_run(); r.text = FOOT_TITLE
            f = r.font; f.name = FONT; f.size = Pt(7.5); f.color.rgb = GREY
            _track(r, 140)
            tb = self._text(s, Inches(12.13), Inches(7.10), Inches(0.6), Inches(0.3))
            p = tb.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
            r = p.add_run(); r.text = str(self.n)
            f = r.font; f.name = FONT; f.size = Pt(9.5); f.bold = True
            f.color.rgb = NAVY
        return s

    def navy_slide(self):
        """Full-bleed navy slide (cover / dividers)."""
        s = self.prs.slides.add_slide(self.blank)
        self.n += 1
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
        bg.line.fill.background(); _flat(bg)
        base = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.92),
                                  SLIDE_W, Inches(0.58))
        base.fill.solid(); base.fill.fore_color.rgb = NAVY_D
        base.line.fill.background(); _flat(base)
        return s

    def motif(self, s, x, y, w=Inches(2.55), h=Inches(0.07)):
        """Four-action ladder strip (exit/hold/quiet/public) — the deck motif."""
        seg = Emu(int(w) // 4)
        for i, col in enumerate((ROSE, SAND, CYAN, TEAL)):
            r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Emu(int(x) + i * int(seg)), y, seg, h)
            r.fill.solid(); r.fill.fore_color.rgb = col
            r.line.fill.background(); _flat(r)

    def hairline(self, s, x, y, w, color=HAIR, weight=0.75):
        # thin filled rectangle (not a connector: shadow-proof in all renderers)
        ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(weight))
        ln.fill.solid(); ln.fill.fore_color.rgb = color
        ln.line.fill.background()
        _flat(ln)
        return ln

    @staticmethod
    def _text(s, x, y, w, h):
        box = s.shapes.add_textbox(x, y, w, h)
        box.text_frame.word_wrap = True
        return box.text_frame

    # ── content elements ────────────────────────────────────────────────────

    def body(self, s, x, y, w, h, items, size=17, gap=6, color=INK,
             align=PP_ALIGN.LEFT):
        """items: list of str or (str, dict) with bold/color/size/level/space."""
        tf = self._text(s, x, y, w, h)
        first = True
        for it in items:
            txt, opt = (it, {}) if isinstance(it, str) else it
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(opt.get("space", gap))
            p.level = opt.get("level", 0)
            p.alignment = opt.get("align", align)
            for seg_txt, seg_bold, seg_col in _segments(txt, opt):
                r = p.add_run(); r.text = seg_txt
                f = r.font; f.name = FONT
                f.size = Pt(opt.get("size", size))
                f.bold = seg_bold
                f.color.rgb = seg_col if seg_col else opt.get("color", color)
        return tf

    def card(self, s, x, y, w, h, accent=None, fill=WHITE):
        """White card with hairline border and optional colored top edge."""
        shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = HAIR; shp.line.width = Pt(0.75)
        _flat(shp)
        if accent is not None:
            bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.045))
            bar.fill.solid(); bar.fill.fore_color.rgb = accent
            bar.line.fill.background(); _flat(bar)
        return shp

    def panel(self, s, x, y, w, h, fill=FAINT, line=None):
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shp.adjustments[0] = 0.035
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line; shp.line.width = Pt(1.0)
        _flat(shp)
        return shp

    def label(self, s, x, y, w, text, color=GREY, size=9.5, spc=160,
              align=PP_ALIGN.LEFT):
        """Tracked small-caps label."""
        tf = self._text(s, x, y, w, Inches(0.3))
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text.upper()
        f = r.font; f.name = FONT; f.size = Pt(size); f.bold = True
        f.color.rgb = color
        _track(r, spc)
        return tf

    def disc(self, s, x, y, text, fill=NAVY, d=Inches(0.52), size=18,
             color=WHITE):
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        shp.line.fill.background(); _flat(shp)
        tf = shp.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        f = r.font; f.name = FONT; f.size = Pt(size); f.bold = True
        f.color.rgb = color
        return shp

    def colorbox(self, s, x, y, w, h, fill):
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shp.adjustments[0] = 0.10
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        shp.line.fill.background(); _flat(shp)
        return shp

    def arrow(self, s, x1, y1, x2, y2, color=GREY, w=Pt(1.75)):
        conn = s.shapes.add_connector(2, x1, y1, x2, y2)  # straight
        conn.line.color.rgb = color; conn.line.width = w
        _flat(conn)
        le = conn.line._get_or_add_ln()
        import copy
        from pptx.oxml.ns import qn
        head = le.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med",
                                                "len": "med"})
        le.append(copy.deepcopy(head))
        return conn

    def picture(self, s, name, x, y, w=None, h=None):
        path = os.path.join(ASSETS, f"{name}.png")
        if not os.path.exists(path):
            ph = self.panel(s, x, y, w or Inches(6), h or Inches(4), fill=FAINT)
            tf = ph.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = f"[figure: {name}]"
            r.font.name = FONT; r.font.size = Pt(14); r.font.color.rgb = GREY
            return ph
        return s.shapes.add_picture(path, x, y, width=w, height=h)

    def caption(self, s, x, y, w, text, size=11):
        tf = self._text(s, x, y, w, Inches(0.5))
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = text
        f = r.font; f.name = FONT; f.size = Pt(size); f.color.rgb = GREY

    def eq(self, s, key, pt, x=None, y=None, cx=None, cy=None):
        """Place a typeset equation PNG at the requested point size.

        Position with (x, y) top-left, or center via cx and/or cy.
        Returns (w, h) as Emu for layout chaining.
        """
        path = self.eqs[key]
        px_w, px_h = eq_render.png_size(path)
        scale = pt / eq_render.BASE_PT
        w = Emu(int(px_w / eq_render.DPI * scale * EMU_IN))
        h = Emu(int(px_h / eq_render.DPI * scale * EMU_IN))
        if cx is not None:
            x = Emu(int(cx) - int(w) // 2)
        if cy is not None:
            y = Emu(int(cy) - int(h) // 2)
        s.shapes.add_picture(path, x, y, width=w, height=h)
        return w, h


# ── deck assembly ───────────────────────────────────────────────────────────

def build() -> None:
    d = Deck()
    f1 = fact1_numbers()
    f2 = fact2_numbers()

    # 1 ── title (full-bleed navy) ───────────────────────────────────────────
    s = d.navy_slide()
    d.motif(s, Inches(0.98), Inches(1.70))
    tf = d._text(s, Inches(0.95), Inches(1.96), Inches(11.8), Inches(1.9))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Liquidity, Activism Disclosure,"
    r.font.name = FONT; r.font.size = Pt(41); r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    r = p2.add_run(); r.text = "and Takeover Premia"
    r.font.name = FONT; r.font.size = Pt(41); r.font.bold = True
    r.font.color.rgb = WHITE
    p3 = tf.add_paragraph(); p3.space_before = Pt(18)
    r = p3.add_run()
    r.text = "When can the market see an activist coming — and who pockets the premium?"
    r.font.name = FONT; r.font.size = Pt(17); r.font.color.rgb = ICE
    tf = d._text(s, Inches(0.95), Inches(5.55), Inches(8), Inches(1.0))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Austin Li"
    r.font.name = FONT; r.font.size = Pt(16); r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.space_before = Pt(3)
    r = p2.add_run(); r.text = "University College London"
    r.font.name = FONT; r.font.size = Pt(13); r.font.color.rgb = ICE
    d.label(s, Inches(9.05), Inches(7.04), Inches(3.7),
            "Seminar deck · 40 minutes", color=ICE, size=8.5,
            align=PP_ALIGN.RIGHT)

    # 2 ── agenda ────────────────────────────────────────────────────────────
    s = d.slide("", "Agenda")
    rows = [
        ("01", "Motivation",
         "Activist returns ride the takeover channel; two live policy levers"),
        ("02", "Framework",
         "One trading round, four plays, and a bidder who reads the price"),
        ("03", "Results",
         "The liquidity sweet spot, the premium wedge, the transparency paradox"),
        ("04", "Evidence",
         "The 2024 five-business-day rule as a natural experiment"),
        ("05", "Implications",
         "Policy, investors, and the road ahead"),
    ]
    y = Inches(1.98)
    for num, head, desc in rows:
        tb = d._text(s, Inches(0.85), y, Inches(0.85), Inches(0.5))
        p = tb.paragraphs[0]
        r = p.add_run(); r.text = num
        r.font.name = FONT; r.font.size = Pt(20); r.font.bold = True
        r.font.color.rgb = NAVY
        d.body(s, Inches(1.85), y + Inches(0.02), Inches(3.3), Inches(0.5),
               [(head, {"bold": True, "size": 17.5})])
        d.body(s, Inches(5.35), y + Inches(0.07), Inches(7.2), Inches(0.5),
               [(desc, {"size": 13.5, "color": GREY})])
        y += Inches(0.62)
        d.hairline(s, Inches(0.85), y, Inches(11.65))
        y += Inches(0.30)
    d.body(s, Inches(0.85), y + Inches(0.05), Inches(11.5), Inches(0.5), [
        ("Backup: calibration, equilibrium tables, posterior algebra, "
         "certification machinery, data & methods.",
         {"size": 12.5, "color": GREY}),
    ])

    # 3 ── the takeaway up front ────────────────────────────────────────────
    s = d.slide("Executive summary", "Three things to take away")
    claims = [
        ("1", "Liquidity has a sweet spot",
         "Takeover gains for minority shareholders are **hump-shaped** in market "
         "liquidity: noise lets bidders in, but too much noise destroys the "
         "market's ability to price activism.", BLUE),
        ("2", "Who pockets the premium is tender mechanics",
         "The activist premium is pinned by a free-rider tender game: "
         "appropriability λ = 1 − q(1−γ)ψ. Low λ explains why estimated "
         "premia can **fall** when activists show up.", TEAL),
        ("3", "Disclosure law is a market-design lever",
         "Stake-triggered disclosure attenuates the liquidity effect — and "
         "stricter rules can **backfire** by killing quiet engagement. The 2024 "
         "SEC acceleration moved exactly this margin.", ROSE),
    ]
    x = MARGIN
    for num, head, body, col in claims:
        d.card(s, x, Inches(1.95), Inches(3.88), Inches(4.55), accent=col)
        d.disc(s, x + Inches(0.28), Inches(2.25), num, fill=col)
        d.body(s, x + Inches(0.28), Inches(2.98), Inches(3.32), Inches(1.1),
               [(head, {"bold": True, "size": 16})])
        d.body(s, x + Inches(0.28), Inches(4.05), Inches(3.32), Inches(2.3),
               [(body, {"size": 13})], color=INK)
        x += Inches(4.13)

    # 4 ── why care ─────────────────────────────────────────────────────────
    s = d.slide("01 · Motivation", "Activist returns ride the takeover channel")
    d.card(s, MARGIN, Inches(1.95), Inches(5.65), Inches(4.55), accent=NAVY)
    d.label(s, Inches(0.90), Inches(2.25), Inches(5.0),
            "Share of activist returns tied to M&A")
    tf = d._text(s, Inches(0.88), Inches(2.55), Inches(5.1), Inches(1.3))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "~65%"
    r.font.name = FONT; r.font.size = Pt(58); r.font.bold = True
    r.font.color.rgb = NAVY
    d.body(s, Inches(0.90), Inches(3.85), Inches(5.05), Inches(2.5), [
        ("of activist hedge-fund returns are tied to M&A outcomes "
         "(Brav–Jiang–Partnoy–Thomas 2008).", {"size": 14.5, "space": 8}),
        ("Targets are acquired at significant premia; activists profit "
         "primarily through the takeover exit (Greenwood–Schor 2009).",
         {"size": 14.5}),
    ])
    d.body(s, Inches(6.75), Inches(2.05), Inches(6.0), Inches(4.6), [
        ("So the question that prices an activist position is:",
         {"size": 16.5, "space": 10}),
        ("does a takeover arrive, and at what premium?",
         {"bold": True, "size": 21, "space": 14}),
        ("Both depend on what the market can infer about engagement "
         "before anything is announced — i.e., on **trading conditions** "
         "and **disclosure rules**.", {"size": 16.5, "space": 10}),
        ("Informed trading moves prices ahead of 13D filings "
         "(Collin-Dufresne–Fos 2015): the market is watching order flow.",
         {"size": 14.5, "color": GREY}),
    ])

    # 5 ── the two levers ───────────────────────────────────────────────────
    s = d.slide("01 · Motivation", "Two policy levers, one information problem")
    d.card(s, MARGIN, Inches(1.95), Inches(5.89), Inches(4.15), accent=NAVY)
    d.label(s, Inches(0.88), Inches(2.25), Inches(5.4),
            "Liquidity — can you trade unseen?", color=NAVY, size=11.5,
            spc=110)
    d.body(s, Inches(0.88), Inches(2.75), Inches(5.35), Inches(3.4), [
        ("Noise trading is **camouflage**: it lets a blockholder accumulate "
         "or exit without moving the price.", {"size": 15, "space": 8}),
        ("But camouflage cuts both ways — it also stops the market from "
         "**seeing value being created**.", {"size": 15, "space": 8}),
        ("Tick sizes, dark pools, retail flow, index churn: all shift this "
         "dial. Regulators move it constantly.", {"size": 13.5, "color": GREY}),
    ])
    d.card(s, Inches(6.84), Inches(1.95), Inches(5.89), Inches(4.15), accent=TEAL)
    d.label(s, Inches(7.12), Inches(2.25), Inches(5.4),
            "Disclosure — when must you show your hand?", color=TEAL,
            size=11.5, spc=110)
    d.body(s, Inches(7.12), Inches(2.75), Inches(5.35), Inches(3.4), [
        ("Schedule 13D: crossing **5%** with intent forces a public filing — "
         "engagement becomes observable.", {"size": 15, "space": 8}),
        ("Below the trigger, engagement can only be **inferred** from order "
         "flow. The rule partitions the information regime.", {"size": 15,
         "space": 8}),
        ("US window: 10 calendar days → **5 business days** (Feb 2024). "
         "UK: 3% / 2 days. The lever is live.", {"size": 13.5, "color": GREY}),
    ])

    # 6 ── 2024 timeline ────────────────────────────────────────────────────
    s = d.slide("01 · Motivation", "2024: the SEC halved the clock")
    y = Inches(3.0)
    d.arrow(s, Inches(0.70), y, Inches(12.60), y, color=GREY, w=Pt(2.0))
    marks = [
        (Inches(1.6), "1968", "Williams Act:\n13D regime born", HAIR),
        (Inches(4.6), "2011–22", "Petitions & debate:\n10-day gap criticized", HAIR),
        (Inches(7.6), "Oct 2023", "Rule 33-11253\nadopted", CYAN),
        (Inches(10.6), "Feb 5, 2024", "5 business days\nin force", NAVY),
    ]
    for x, when, what, accent in marks:
        tick = s.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.075),
                                  y - Inches(0.075), Inches(0.15), Inches(0.15))
        tick.fill.solid(); tick.fill.fore_color.rgb = NAVY
        tick.line.color.rgb = WHITE; tick.line.width = Pt(1.5)
        _flat(tick)
        d.body(s, x - Inches(1.05), y - Inches(0.72), Inches(2.1), Inches(0.4),
               [(when, {"bold": True, "size": 14, "color": NAVY,
                        "align": PP_ALIGN.CENTER})])
        d.card(s, x - Inches(1.05), y + Inches(0.28), Inches(2.1), Inches(1.05),
               accent=accent)
        d.body(s, x - Inches(0.95), y + Inches(0.42), Inches(1.9), Inches(0.85),
               [(what, {"size": 11.5, "align": PP_ALIGN.CENTER})])
    d.body(s, Inches(0.9), Inches(5.15), Inches(11.6), Inches(1.6), [
        ("A regulator just compressed the **pre-disclosure accumulation "
         "window** — the exact margin this framework prices.",
         {"size": 17.5, "space": 8}),
        ("That makes the theory testable now: filings, returns, and "
         "liquidity around 2024-02-05 are the natural experiment.",
         {"size": 15, "color": GREY}),
    ])

    # 7 ── the question ─────────────────────────────────────────────────────
    s = d.slide("01 · Motivation", "", footer=True)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.05), Inches(2.12),
                             Inches(0.07), Inches(2.30))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background(); _flat(bar)
    d.label(s, Inches(1.45), Inches(2.08), Inches(6), "The question")
    tf = d._text(s, Inches(1.42), Inches(2.45), Inches(10.4), Inches(1.9))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = ("How do market liquidity and disclosure rules shape the takeover "
              "premia that minority shareholders receive?")
    r.font.name = FONT; r.font.size = Pt(27); r.font.bold = True
    r.font.color.rgb = INK
    d.body(s, Inches(1.45), Inches(4.85), Inches(10.2), Inches(1.2), [
        ("A single equilibrium framework: a blockholder who can exit, hold, or "
         "engage — quietly or publicly — a market pricing order flow, and a "
         "bidder reading both.", {"size": 15.5, "color": GREY}),
    ])

    # 8 ── positioning vs the literature ────────────────────────────────────
    s = d.slide("01 · Positioning", "Three literatures, one missing piece")
    cols = [
        ("Exit vs voice", "Blockholder governance: liquidity both "
         "disciplines (cheap accumulation) and tempts exit.\n"
         "Hirschman 1970; Maug 1998; Edmans 2009", ROSE),
        ("Informed trading", "Microstructure: order flow moves prices; "
         "noise is camouflage.\nKyle 1985; Glosten–Milgrom 1985", CYAN),
        ("Takeover feedback", "Prices feed back into real decisions — "
         "including bids.\nGrossman–Hart 1980; Edmans–Goldstein–Jiang 2015",
         TEAL),
    ]
    x = MARGIN
    for head, body_txt, col in cols:
        d.card(s, x, Inches(1.95), Inches(3.88), Inches(2.95), accent=col)
        d.label(s, x + Inches(0.26), Inches(2.22), Inches(3.4), head,
                color=INK, size=12)
        main, _, cites = body_txt.partition("\n")
        d.body(s, x + Inches(0.26), Inches(2.72), Inches(3.38), Inches(2.0), [
            (main, {"size": 13.5, "space": 8}),
            (cites, {"size": 11, "color": GREY}),
        ])
        x += Inches(4.13)
    d.panel(s, MARGIN, Inches(5.20), CONTENT_W, Inches(1.5))
    d.body(s, Inches(0.90), Inches(5.40), Inches(11.6), Inches(1.2), [
        ("Closest trio: Ordóñez-Calafí & Bernhardt 2022 (threshold design), "
         "Corum & Levit 2019 (activists as bidder catalysts), Cetemen et al. "
         "2026 (dynamic trade timing).", {"size": 13, "color": GREY,
                                          "space": 6}),
        ("**This paper:** stake-triggered disclosure partitions order-flow "
         "inference into disclosed and nondisclosed branches — and that "
         "inference feeds endogenous bidder entry.", {"size": 15}),
    ])

    # 9 ── framework in one picture ─────────────────────────────────────────
    s = d.slide("02 · Framework", "The model in one picture")
    stages = [
        ("Nature", "Value v drawn; blockholder sees private signal s", HAIR),
        ("Trading", "Chooses action; noise traders mix in; market prices "
         "order flow", NAVY),
        ("Entry", "Bidder reads price & disclosure, decides whether to bid",
         TEAL),
        ("Payoffs", "Takeover premium or standalone value", HAIR),
    ]
    x = MARGIN
    for i, (head, body, accent) in enumerate(stages):
        d.card(s, x, Inches(2.15), Inches(2.62), Inches(2.25), accent=accent)
        d.label(s, x + Inches(0.22), Inches(2.42), Inches(2.2), head,
                color=NAVY if accent is not HAIR else INK, size=12)
        d.body(s, x + Inches(0.22), Inches(2.92), Inches(2.2), Inches(1.4),
               [(body, {"size": 12.5})])
        if i < 3:
            d.arrow(s, x + Inches(2.70), Inches(3.27), x + Inches(3.10),
                    Inches(3.27), color=GREY)
        x += Inches(3.17)
    d.body(s, MARGIN, Inches(5.05), Inches(12), Inches(1.7), [
        ("One round of trading, one potential bidder — the cleanest lab for the "
         "question. κ = share of noise trades (the liquidity dial).",
         {"size": 15.5, "space": 8}),
        ("Crucially, the bidder **conditions on the price**: information flows "
         "from trading to takeover entry. That feedback loop is the engine.",
         {"size": 15.5}),
    ])

    # 10 ── the four plays ──────────────────────────────────────────────────
    s = d.slide("02 · Framework", "The blockholder's four plays")
    plays = [
        ("Exit", "Sell the stake", "Weakest signals: cash out before bad news",
         ROSE),
        ("Hold", "Stay passive", "Middling signals: not worth the engagement "
         "cost", SAND),
        ("Quiet voice", "Engage privately", "Good signals: improve the firm "
         "below the radar", CYAN),
        ("Public voice", "Buy + engage + file", "Best signals: scale up and "
         "trigger disclosure", TEAL),
    ]
    x = MARGIN
    for head, sub, body, col in plays:
        d.card(s, x, Inches(1.95), Inches(2.92), Inches(3.15), accent=col)
        d.label(s, x + Inches(0.24), Inches(2.22), Inches(2.5), head,
                color=INK, size=12)
        d.body(s, x + Inches(0.24), Inches(2.72), Inches(2.45), Inches(0.55),
               [(sub, {"bold": True, "size": 14})])
        d.body(s, x + Inches(0.24), Inches(3.30), Inches(2.45), Inches(1.55),
               [(body, {"size": 12.5})])
        x += Inches(3.07)
    d.arrow(s, Inches(1.2), Inches(5.62), Inches(12.15), Inches(5.62),
            color=GREY, w=Pt(1.75))
    tb = d._text(s, Inches(2.2), Inches(5.82), Inches(2.9), Inches(0.4))
    p = tb.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = "private signal s  →  cutoff ladder"
    r.font.name = FONT; r.font.size = Pt(13.5); r.font.color.rgb = GREY
    d.eq(s, "ladder", 15, cx=Inches(6.67), cy=Inches(5.94))
    tb = d._text(s, Inches(8.2), Inches(5.82), Inches(4.0), Inches(0.4))
    p = tb.paragraphs[0]
    r = p.add_run(); r.text = "(a theorem, not an assumption)"
    r.font.name = FONT; r.font.size = Pt(13.5); r.font.color.rgb = GREY

    # 11 ── the feedback loop ───────────────────────────────────────────────
    s = d.slide("02 · Framework", "Why liquidity matters twice")
    nodes = [
        ("Order flow", "node_X", Inches(0.60), NAVY, ICE, 18),
        ("Market belief", "node_belief", Inches(3.83), CYAN, INK, 15),
        ("Price", "node_price", Inches(7.05), SAND, INK, 15),
        ("Bid entry", "node_p", Inches(10.28), TEAL, WHITE, 18),
    ]
    for txt, key, x, fill, labcol, pt in nodes:
        d.colorbox(s, x, Inches(2.75), Inches(2.45), Inches(1.30), fill)
        d.label(s, x, Inches(2.92), Inches(2.45), txt, color=labcol,
                size=10.5, align=PP_ALIGN.CENTER)
        d.eq(s, key, pt, cx=Emu(int(x) + int(Inches(1.225))), cy=Inches(3.68))
        if x < Inches(10):
            d.arrow(s, x + Inches(2.52), Inches(3.40), x + Inches(3.76),
                    Inches(3.40), color=GREY)
    d.body(s, MARGIN, Inches(4.60), Inches(12.1), Inches(2.3), [
        ("**Noise as cover (+):** more noise → lower prices on average → "
         "takeovers get cheaper → more bids arrive.", {"size": 16,
         "space": 8}),
        ("**Noise as fog (−):** more noise → order flow says less about "
         "engagement → the activism premium the market can price **erodes**.",
         {"size": 16, "space": 8}),
        ("Disclosure (D = 1) bypasses the fog entirely — engagement is "
         "observed, not inferred. That asymmetry drives the policy results.",
         {"size": 14.5, "color": GREY}),
    ])

    # 12 ── two regimes ─────────────────────────────────────────────────────
    s = d.slide("02 · Framework", "What the market sees: two information regimes")
    d.card(s, MARGIN, Inches(1.95), Inches(5.89), Inches(3.55), accent=CYAN)
    d.label(s, Inches(0.88), Inches(2.25), Inches(5.4),
            "Quiet (D = 0) — inference", color=INK, size=12.5)
    d.body(s, Inches(0.88), Inches(2.80), Inches(5.35), Inches(3.0), [
        ("Engagement is **guessed from order flow**.", {"size": 15.5,
         "space": 8}),
        ("Liquidity κ governs how good the guess is.", {"size": 15.5,
         "space": 8}),
        ("Hold and quiet voice look identical in the tape — by design.",
         {"size": 14, "color": GREY}),
    ])
    d.card(s, Inches(6.84), Inches(1.95), Inches(5.89), Inches(3.55), accent=TEAL)
    d.label(s, Inches(7.12), Inches(2.25), Inches(5.4),
            "Public (D = 1) — observation", color=INK, size=12.5)
    d.body(s, Inches(7.12), Inches(2.80), Inches(5.35), Inches(3.0), [
        ("The 13D filing makes engagement **public**.", {"size": 15.5,
         "space": 8}),
        ("Prices stop depending on κ at all.", {"size": 15.5, "space": 8}),
        ("Disclosure is an information-regime switch, not just paperwork.",
         {"size": 14, "color": GREY}),
    ])
    d.body(s, MARGIN, Inches(5.90), Inches(12.0), Inches(0.6), [
        ("The disclosure rule decides **how much of the market lives in each "
         "regime** — that is what makes it a design object.", {"size": 15.5}),
    ])

    # 13 ── result 1: the hump ──────────────────────────────────────────────
    s = d.slide("03 · Results · 1", "Liquidity has a sweet spot")
    d.picture(s, "nonmonotone", MARGIN, Inches(1.95), w=Inches(6.2))
    d.body(s, Inches(7.15), Inches(2.0), Inches(5.55), Inches(2.9), [
        ("Expected takeover gains to minority shareholders are "
         "**hump-shaped** in liquidity.", {"size": 16.5, "space": 10}),
        ("Too little noise: bidders stay away — prices are too informative, "
         "deals too expensive.", {"size": 15}),
        ("Too much noise: bids come, but the **activism premium** the market "
         "can price has evaporated.", {"size": 15}),
    ])
    d.card(s, Inches(7.15), Inches(5.05), Inches(5.55), Inches(1.30),
           accent=NAVY)
    d.label(s, Inches(7.43), Inches(5.28), Inches(2.4), "The peak")
    d.eq(s, "kdagger", 22, x=Inches(7.43), cy=Inches(5.92))
    d.body(s, Inches(9.55), Inches(5.62), Inches(3.05), Inches(0.7), [
        ("cover and visibility in balance", {"size": 12.5, "color": GREY}),
    ])
    d.caption(s, MARGIN, Inches(6.55), Inches(7),
              "Minority takeover gains Δ-min vs noise-trading intensity κ; baseline calibration.")

    # 14 ── why: decomposition ──────────────────────────────────────────────
    s = d.slide("03 · Results · 1", "Why: cover vs camouflage, in one split")
    d.picture(s, "decomposition", MARGIN, Inches(1.95), w=Inches(6.2))
    d.body(s, Inches(7.15), Inches(2.0), Inches(5.55), Inches(4.4), [
        ("Split the gains into two books:", {"size": 15.5, "space": 10}),
        ("**Base book** — premium that arrives regardless of activism. "
         "Rises with noise: cover brings bids.", {"size": 15, "space": 10}),
        ("**Activism book** — premium tied to the market believing engagement "
         "is happening. Hump-shaped: inference dies with noise.",
         {"size": 15, "space": 10}),
        ("The activism book's peak sets the overall sweet spot.",
         {"size": 15}),
    ])
    d.caption(s, MARGIN, Inches(6.55), Inches(7),
              "Decomposition of Δ-min into base and activism components across κ.")

    # 15 ── result 2: the wedge ─────────────────────────────────────────────
    s = d.slide("03 · Results · 2", "Who captures the prize: tender-game arithmetic")
    d.picture(s, "wedge", MARGIN, Inches(1.95), w=Inches(6.85))
    bx, bw = Inches(7.70), Inches(5.03)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, Inches(1.95), bw,
                              Inches(1.52))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.line.fill.background(); _flat(band)
    d.eq(s, "lam_white", 23, cx=Emu(int(bx) + int(bw) // 2), cy=Inches(2.52))
    d.body(s, bx, Inches(2.98), bw, Inches(0.45), [
        ("how much of the improvement the activist's side keeps",
         {"size": 11, "color": ICE, "align": PP_ALIGN.CENTER}),
    ])
    legend = [
        ("sym_q", "— chance a rival raider shows up"),
        ("sym_gamma", "— does the improvement survive a change of control?"),
        ("sym_psi", "— can the bloc's stake block the raid?"),
    ]
    y = Inches(3.76)
    for key, txt in legend:
        d.eq(s, key, 14, x=Inches(7.78), cy=Emu(int(y) + int(Inches(0.16))))
        d.body(s, Inches(8.18), y, Inches(4.55), Inches(0.5),
               [(txt, {"size": 13})])
        y += Inches(0.56)
    d.label(s, Inches(7.78), Inches(5.50), Inches(3.5), "Premium wedge")
    d.eq(s, "wedge", 14, x=Inches(7.78), y=Inches(5.80))
    d.body(s, Inches(7.78), Inches(6.18), Inches(4.9), Inches(0.5), [
        ("derived from a free-rider tender game — not assumed.",
         {"size": 12, "color": GREY}),
    ])
    d.caption(s, MARGIN, Inches(6.55), Inches(7),
              "Appropriability λ across game primitives (q, γ, ψ).")

    # 16 ── result 2b: negative premia reconciled ───────────────────────────
    s = d.slide("03 · Results · 2", "The puzzle this solves: premia that fall on arrival")
    d.body(s, MARGIN, Inches(1.95), Inches(5.95), Inches(4.4), [
        ("Structural estimates (Celentano–Levine 2025) find bid premia "
         "**13.7% lower** with an activist in the deal.", {"size": 16.5,
         "space": 10}),
        ("Contradiction with 'activism creates value'? No.", {"size": 15.5,
         "space": 8}),
        ("The measured premium divides the offer by a market price that "
         "**already banks** the activist's work (AFS 2022: 13D returns are "
         "~3/4 anticipated treatment).", {"size": 15.5, "space": 8}),
        ("When appropriability λ is low, the denominator outruns the "
         "numerator: measured premia fall even though the true wedge "
         "m₁ ≥ m₀ never flips.", {"size": 15.5}),
    ])
    d.card(s, Inches(7.05), Inches(2.0), Inches(5.67), Inches(4.25),
           accent=TEAL)
    d.body(s, Inches(7.35), Inches(2.28), Inches(5.05), Inches(3.8), [
        ("What our wedge adds vs the board-bargaining story:",
         {"bold": True, "size": 14.5, "space": 10}),
        ("Their effect moves only with activist **presence**.", {"size": 14,
         "space": 8}),
        ("Ours is indexed by **observables**: acquirer type (γ), rival "
         "pressure (q), stake pivotality (ψ), liquidity (κ).", {"size": 14,
         "space": 8}),
        ("→ It predicts **where the sign flips**: portable improvements, "
         "pivotal blocs, cold fringe markets.", {"size": 14, "space": 8}),
        ("Different cross-sections, testable against each other.",
         {"size": 13, "color": GREY}),
    ])

    # 17 ── result 3: when a theorem ────────────────────────────────────────
    s = d.slide("03 · Results · 3", "How robust is the sweet spot? Now a theorem — with edges")
    d.picture(s, "ge_decomp", MARGIN, Inches(1.95), w=Inches(6.85))
    d.body(s, Inches(7.70), Inches(2.0), Inches(5.03), Inches(2.7), [
        ("On a **certified range** of the liquidity dial, the hump is a "
         "theorem: equilibrium feedback provably cannot overturn it.",
         {"size": 15, "space": 10}),
        ("Outside it, the warning is real: when bidders barely react to "
         "premia, the hump **inverts into a trough** — certified, not "
         "hypothetical.", {"size": 15}),
    ])
    d.label(s, Inches(7.70), Inches(4.78), Inches(2.5), "Certified range")
    d.eq(s, "range_cert", 15, x=Inches(7.70), cy=Inches(5.24))
    d.label(s, Inches(10.45), Inches(4.78), Inches(2.3), "Inversion case")
    d.eq(s, "sigxi", 15, x=Inches(10.45), cy=Inches(5.24))
    d.body(s, Inches(7.70), Inches(5.70), Inches(5.03), Inches(1.0), [
        ("Practical read: the liquidity sweet spot is strongest where "
         "takeover entry is most sensitive to expected resistance.",
         {"size": 13.5, "color": GREY}),
    ])
    d.caption(s, MARGIN, Inches(6.55), Inches(7),
              "GE channel decomposition; certified region and the σ_ξ = 0.60 inversion.")

    # 18 ── disclosure attenuates ───────────────────────────────────────────
    s = d.slide("03 · Results · 4", "Disclosure turns the liquidity dial down")
    d.picture(s, "slopes", MARGIN, Inches(1.95), w=Inches(6.85))
    d.body(s, Inches(7.70), Inches(2.0), Inches(5.03), Inches(4.4), [
        ("Disclosed states don't depend on liquidity — so the more "
         "engagement happens in public, the **flatter** the "
         "liquidity–premium relationship.", {"size": 15.5, "space": 12}),
        ("Disclosure and liquidity are **substitute** information channels.",
         {"size": 15.5, "space": 12}),
        ("Cross-market prediction: tighter regimes (UK 3%/2-day) should show "
         "**muted** liquidity sensitivity vs the US.", {"size": 14.5}),
    ])
    d.caption(s, MARGIN, Inches(6.55), Inches(7),
              "Sensitivity of activism gains to κ under baseline vs stricter disclosure.")

    # 19 ── the policy paradox ──────────────────────────────────────────────
    s = d.slide("03 · Results · 4", "The transparency paradox")
    d.card(s, MARGIN, Inches(2.0), Inches(5.89), Inches(2.80), accent=TEAL)
    d.label(s, Inches(0.88), Inches(2.28), Inches(5.3), "Transparency (+)",
            color=TEAL, size=12.5)
    d.body(s, Inches(0.88), Inches(2.80), Inches(5.35), Inches(2.2), [
        ("Stricter disclosure reveals engagement → market prices activism "
         "→ minorities share the gains.", {"size": 15, "space": 8}),
        ("Dominates when engagement is cheap and liquidity moderate.",
         {"size": 13.5, "color": GREY}),
    ])
    d.card(s, Inches(6.84), Inches(2.0), Inches(5.89), Inches(2.80), accent=ROSE)
    d.label(s, Inches(7.12), Inches(2.28), Inches(5.3), "Deterrence (−)",
            color=ROSE, size=12.5)
    d.body(s, Inches(7.12), Inches(2.80), Inches(5.35), Inches(2.2), [
        ("Forced visibility kills **quiet voice**: the blockholder exits "
         "rather than tip a bidder-deterring hand.", {"size": 15, "space": 8}),
        ("Dominates when engagement is costly — minorities end up worse off.",
         {"size": 13.5, "color": GREY}),
    ])
    d.panel(s, MARGIN, Inches(5.30), CONTENT_W, Inches(1.05))
    d.body(s, Inches(0.90), Inches(5.56), Inches(11.5), Inches(0.8), [
        ("Disclosure reform that ignores deterrence can **destroy the channel "
         "it is trying to illuminate** — the model says when.",
         {"size": 16}),
    ])

    # 20 ── evidence 1 ──────────────────────────────────────────────────────
    s = d.slide("04 · Evidence · 1", "The clock bit: filings sped up immediately")
    d.picture(s, "fact1", MARGIN, Inches(1.95), w=Inches(6.85))
    d.card(s, Inches(7.70), Inches(1.95), Inches(5.03), Inches(2.55),
           accent=NAVY)
    d.label(s, Inches(7.98), Inches(2.20), Inches(4.5), "Median filing delay")
    tf = d._text(s, Inches(7.98), Inches(2.50), Inches(4.5), Inches(0.75))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"{f1['pre_med']} → {f1['post_med']} days"
    r.font.name = FONT; r.font.size = Pt(30); r.font.bold = True
    r.font.color.rgb = NAVY
    d.label(s, Inches(7.98), Inches(3.42), Inches(4.5),
            "Filed within 5 business days")
    tf = d._text(s, Inches(7.98), Inches(3.72), Inches(4.5), Inches(0.6))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"{f1['pre_5bd']} → {f1['post_5bd']}"
    r.font.name = FONT; r.font.size = Pt(24); r.font.bold = True
    r.font.color.rgb = INK
    d.body(s, Inches(7.70), Inches(4.80), Inches(5.03), Inches(1.7), [
        ("Universe of original 13Ds on EDGAR, windows straddling the "
         "2024-02-05 compliance date.", {"size": 13.5, "color": GREY,
         "space": 8}),
        ("The accumulation window — the model's quiet-trading cover — "
         "shrank by force of law.", {"size": 14}),
    ])

    # 21 ── evidence 2 ──────────────────────────────────────────────────────
    s = d.slide("04 · Evidence · 2", "The market noticed: announcement returns moved")
    d.picture(s, "fact2", MARGIN, Inches(1.95), w=Inches(6.85))
    d.body(s, Inches(7.70), Inches(2.0), Inches(5.03), Inches(2.2), [
        ("Market-model CARs around 1,513 original 13D filings, 2022–2025; "
         "Post = filed under the 5-day rule.", {"size": 13.5, "space": 8}),
        ("Mean CAR[−1,+1] **doubles** post-rule (0.85% → 1.71%; sign as "
         "predicted, not significant). Surprises concentrate in **liquid** "
         "names — the camouflage margin is first-order. The Post×illiquidity "
         "interaction is a precise null.", {"size": 14}),
    ])
    d.card(s, Inches(7.70), Inches(4.42), Inches(5.03), Inches(2.0),
           accent=NAVY)
    d.body(s, Inches(7.98), Inches(4.62), Inches(4.5), Inches(1.7), [
        (f"β (Post): {f2['beta']}   ({f2['beta_t']})", {"size": 14.5,
         "bold": True, "space": 7}),
        (f"ln Amihud: {f2['amihud']}   ({f2['amihud_t']})",
         {"size": 14.5, "bold": True, "space": 7}),
        (f"δ (Post × ln Amihud): {f2['delta']}   ({f2['delta_t']})",
         {"size": 14.5, "bold": True, "space": 7}),
        (f"n = {f2['n']} events · two-way clustered SEs",
         {"size": 11.5, "color": GREY}),
    ])

    # 22 ── policy meaning ──────────────────────────────────────────────────
    s = d.slide("05 · Implications", "For policy: liquidity regulation is governance policy")
    points = [
        ("Anything that moves noise trading — tick size regimes, dark-pool "
         "rules, PFOF, index flows — **also moves takeover discipline**. The "
         "two rulebooks are one.", NAVY),
        ("Disclosure thresholds and windows are not transparency knobs alone: "
         "they reallocate engagement between quiet and public — and can "
         "**deter it outright**.", TEAL),
        ("The 2024 acceleration traded camouflage for visibility. The model "
         "prices that trade; the early data move the predicted way.", ROSE),
    ]
    y = Inches(2.05)
    for i, (txt, col) in enumerate(points, 1):
        d.disc(s, Inches(0.85), y, str(i), fill=col, d=Inches(0.46), size=15)
        d.body(s, Inches(1.65), y - Inches(0.04), Inches(11.0), Inches(1.2),
               [(txt, {"size": 17})])
        y += Inches(1.26)
    d.panel(s, Inches(0.85), y + Inches(0.12), Inches(11.65), Inches(0.80))
    d.body(s, Inches(1.10), y + Inches(0.30), Inches(11.2), Inches(0.6), [
        ("Welfare caveat: the κ that maximizes minority gains is **not** the "
         "social optimum — extraction deters some efficient deals.",
         {"size": 14, "color": GREY}),
    ])

    # 23 ── investor meaning ────────────────────────────────────────────────
    s = d.slide("05 · Implications", "For investors: where activism pays")
    rows = [
        ("Screen on liquidity", "Activist targets in **mid-liquidity** names "
         "carry the largest expected takeover kicker — both tails are weak.",
         BLUE),
        ("Read the wedge", "Premium capture is largest with portable "
         "improvements (high γ), pivotal blocs (ψ), and few rival raiders "
         "(low q).", TEAL),
        ("Mind the regime", "Post-2024, filing-day surprises are larger on "
         "average — and they are largest in liquid names, where quiet "
         "accumulation hides best.", ROSE),
    ]
    y = Inches(2.05)
    for head, body, col in rows:
        sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85),
                                y + Inches(0.10), Inches(0.09), Inches(0.85))
        sq.fill.solid(); sq.fill.fore_color.rgb = col
        sq.line.fill.background(); _flat(sq)
        d.label(s, Inches(1.20), y + Inches(0.06), Inches(3.2), head,
                color=INK, size=13)
        d.body(s, Inches(4.55), y + Inches(0.02), Inches(8.1), Inches(1.2),
               [(body, {"size": 15.5})])
        y += Inches(1.38)
        if y < Inches(6):
            d.hairline(s, Inches(0.85), y - Inches(0.22), Inches(11.65))
    d.body(s, Inches(0.85), Inches(6.30), Inches(11.5), Inches(0.5), [
        ("All three are cross-sectional, sign-testable statements — not "
         "calibration folklore.", {"size": 13.5, "color": GREY}),
    ])

    # 24 ── where this goes ─────────────────────────────────────────────────
    s = d.slide("05 · Roadmap", "Where this goes next")
    steps = [
        ("Now", "Theory + certified numerics + two facts from the 2024 "
         "natural experiment", BLUE),
        ("Next", "Full event-study cross-sections: premia, bid hazards, and "
         "liquidity interactions on the 13D universe", TEAL),
        ("Then", "Structural leg: estimate the engagement-cost distribution "
         "from 13D XML microdata — the white space the structural trio "
         "left open", ROSE),
    ]
    x = MARGIN
    for when, what, col in steps:
        d.card(s, x, Inches(2.10), Inches(3.88), Inches(3.35), accent=col)
        d.label(s, x + Inches(0.26), Inches(2.38), Inches(2.0), when,
                color=col if col is not BLUE else NAVY, size=13)
        d.body(s, x + Inches(0.26), Inches(2.90), Inches(3.35), Inches(2.3),
               [(what, {"size": 14.5})])
        x += Inches(4.13)
    d.body(s, MARGIN, Inches(5.85), Inches(11.9), Inches(0.8), [
        ("The framework is built so each leg sharpens the same objects: "
         "κ, the disclosure rule, and λ.", {"size": 14.5, "color": GREY}),
    ])

    # 25 ── takeaways ───────────────────────────────────────────────────────
    s = d.slide("Close", "Three things to remember")
    msgs = [
        ("Liquidity has a sweet spot", "minority takeover gains peak at "
         "moderate noise — and we certified when that's a theorem.", BLUE,
         None),
        ("Premia follow tender mechanics", "appropriability prices who keeps "
         "the improvement — and why measured premia can fall.", TEAL, "lam"),
        ("Disclosure is market design", "it substitutes for liquidity in "
         "pricing activism, and stricter isn't always better.", ROSE, None),
    ]
    y = Inches(2.15)
    for i, (head, body, col, eq_key) in enumerate(msgs, 1):
        d.disc(s, Inches(0.85), y, str(i), fill=col)
        tw = Inches(8.85) if eq_key else Inches(11)
        d.body(s, Inches(1.70), y - Inches(0.02), tw, Inches(1.2), [
            (f"**{head}** — {body}", {"size": 18.5}),
        ])
        if eq_key:
            d.eq(s, eq_key, 14, x=Inches(10.80),
                 cy=Emu(int(y) + int(Inches(0.30))))
        y += Inches(1.28)
    d.hairline(s, Inches(0.85), Inches(6.10), Inches(11.65))
    d.body(s, Inches(0.85), Inches(6.28), Inches(11), Inches(0.6), [
        ("Thank you  ·  junyu.li.24@ucl.ac.uk", {"size": 15, "color": GREY}),
    ])

    # ── BACKUPS ─────────────────────────────────────────────────────────────
    s = d.navy_slide()
    d.motif(s, Inches(0.98), Inches(3.02))
    tf = d._text(s, Inches(0.95), Inches(3.25), Inches(11.5), Inches(1.0))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Backup"
    r.font.name = FONT; r.font.size = Pt(38); r.font.bold = True
    r.font.color.rgb = WHITE
    d.body(s, Inches(0.98), Inches(4.30), Inches(11), Inches(0.5), [
        ("Calibration · equilibrium tables · posterior algebra · "
         "certification machinery · tender game · data & methods",
         {"size": 13, "color": ICE}),
    ])

    # B1 calibration
    s = d.slide("Backup", "Baseline calibration")
    cal = [
        ("cal_prior", "Prior mean / value vol", "1.00 / 0.50",
         "normalization; moderate uncertainty"),
        ("cal_eps", "Signal noise", "0.50", "signal weight β = 0.5"),
        ("cal_delta", "Discount", "0.95", "near-term horizon"),
        ("cal_kappa", "Liquidity (baseline)", "0.50", "interior"),
        ("cal_cost", "Engagement cost, slope", "0.12, 0.50",
         "moderate, strong-thesis decay"),
        ("cal_rho", "Success prob / value-add", "0.90 / 0.25",
         "selective activists"),
        ("cal_m", "Premia (base / success)", "0.10 / 0.30",
         "10% base, 30% on success"),
        ("cal_syn", "Synergy / entry cost / bidder vol", "1.44 / 0.15 / 0.40",
         "interior bid probability"),
    ]
    d.label(s, Inches(0.85), Inches(1.95), Inches(1.6), "Symbol")
    d.label(s, Inches(2.55), Inches(1.95), Inches(3.5), "Parameter")
    d.label(s, Inches(6.95), Inches(1.95), Inches(1.8), "Value")
    d.label(s, Inches(8.95), Inches(1.95), Inches(3.5), "Rationale")
    d.hairline(s, Inches(0.85), Inches(2.30), Inches(11.65), color=NAVY,
               weight=1.25)
    y = Inches(2.44)
    for key, name, val, why in cal:
        d.eq(s, key, 13, x=Inches(0.87), cy=Emu(int(y) + int(Inches(0.145))))
        d.body(s, Inches(2.55), y, Inches(4.2), Inches(0.42),
               [(name, {"size": 13.5})])
        d.body(s, Inches(6.95), y, Inches(1.9), Inches(0.42),
               [(val, {"size": 13.5, "bold": True})])
        d.body(s, Inches(8.95), y, Inches(3.6), Inches(0.42),
               [(why, {"size": 12, "color": GREY})])
        y += Inches(0.475)
        d.hairline(s, Inches(0.85), y - Inches(0.065), Inches(11.65))
    yy = Emu(int(y) + int(Inches(0.10)))
    tb = d._text(s, Inches(0.85), yy, Inches(1.8), Inches(0.4))
    p = tb.paragraphs[0]
    r = p.add_run(); r.text = "Regularity (A5):"
    r.font.name = FONT; r.font.size = Pt(12.5); r.font.bold = True
    r.font.color.rgb = INK
    d.eq(s, "a5", 11.5, x=Inches(2.45),
         cy=Emu(int(yy) + int(Inches(0.125))))
    d.body(s, Inches(4.85), yy, Inches(7.5), Inches(0.4), [
        ("— unique pricing fixed point.", {"size": 12.5, "color": GREY}),
    ])

    # B2 baseline equilibrium table
    s = d.slide("Backup", "Baseline equilibrium outcomes")
    xcols = [Inches(0.95), Inches(2.15), Inches(3.75), Inches(5.75),
             Inches(7.75), Inches(9.95)]
    headers = ["th_D", "th_X", "th_P", "th_pi", "th_p", "th_mbar"]
    sub = ["disclosure", "order flow", "price", "belief", "bid prob",
           "premium"]
    for x, key, lab in zip(xcols, headers, sub):
        d.eq(s, key, 15, x=x, cy=Inches(2.13))
        d.body(s, x, Inches(2.30), Inches(1.9), Inches(0.3),
               [(lab, {"size": 10.5, "color": GREY})])
    d.hairline(s, Inches(0.85), Inches(2.66), Inches(11.65), color=NAVY,
               weight=1.25)
    data = [
        ("0", "−2", "0.84", "0.00", "0.81", "0.10"),
        ("0", "−1", "1.06", "0.41", "0.55", "0.17"),
        ("0", "0", "1.23", "0.74", "0.33", "0.23"),
        ("0", "+1", "1.39", "1.00", "0.17", "0.28"),
        ("1", "0/1/2", "1.90", "1.00", "0.01", "0.28"),
    ]
    y = Inches(2.80)
    for row in data:
        for x, cell in zip(xcols, row):
            d.body(s, x, y, Inches(1.9), Inches(0.4), [(cell, {"size": 14})])
        y += Inches(0.50)
        d.hairline(s, Inches(0.85), y - Inches(0.07), Inches(11.65))
    d.body(s, Inches(0.85), Emu(int(y) + int(Inches(0.12))), Inches(11.6),
           Inches(1.4), [
        ("Cutoffs: k₁ = k₀ ≈ 0.82, k_D ≈ 2.26 (hold region collapsed at "
         "baseline).", {"size": 13.5, "space": 6}),
        ("Disclosed branch: flat price, belief 1, bid probability collapses "
         "to 1% — public activism strongly deters bidders (price run-up plus "
         "the full premium m̃ = 0.28).", {"size": 13.5}),
    ])

    # B3 posteriors
    s = d.slide("Backup", "Where κ enters: the posterior algebra")
    d.label(s, Inches(0.85), Inches(1.95), Inches(3.0), "Noise")
    d.eq(s, "noise_cases", 14, x=Inches(0.87), y=Inches(2.28))
    d.label(s, Inches(7.05), Inches(1.95), Inches(4.0),
            "Order flow & disclosure")
    d.eq(s, "flow", 14, x=Inches(7.07), y=Inches(2.45))
    d.hairline(s, Inches(0.85), Inches(3.30), Inches(11.65))
    posts = [
        ("post_p1", "κ-free — X = 1 pins q = 0", GREY),
        ("post_0", "falls in κ — noise dilutes the good signal", ROSE),
        ("post_m1", "rises in κ — sell-side flow gets excused", TEAL),
    ]
    y = Inches(3.55)
    for key, note, col in posts:
        d.eq(s, key, 14, x=Inches(0.87), cy=Emu(int(y) + int(Inches(0.30))))
        d.body(s, Inches(8.35), Emu(int(y) + int(Inches(0.17))), Inches(4.3),
               Inches(0.6), [(note, {"size": 13, "color": col, "bold": True})])
        y += Inches(0.92)
    d.body(s, Inches(0.85), Inches(6.30), Inches(11.65), Inches(0.6), [
        ("Rising κ compresses posteriors toward the prior: the inference "
         "channel of the hump. On D = 1, π ≡ 1 regardless of κ.",
         {"size": 13.5, "color": GREY}),
    ])

    # B4 existence/uniqueness
    s = d.slide("Backup", "Existence and uniqueness")
    d.body(s, Inches(0.85), Inches(2.0), Inches(11.65), Inches(3.6), [
        ("Existence: Brouwer on the compact cutoff simplex — unconditional.",
         {"size": 15.5, "space": 10}),
        ("Pricing layer: A5 makes the price map a contraction — unique P* "
         "per information set:", {"size": 15.5, "space": 6}),
    ])
    d.eq(s, "a5", 16.5, cx=Inches(6.67), y=Inches(3.28))
    d.body(s, Inches(0.85), Inches(3.95), Inches(11.65), Inches(2.6), [
        ("Cutoff layer: contraction verified numerically (A6); multi-start "
         "search from 4+ initial conditions converges to identical cutoffs; "
         "modulus L ≤ 0.836 along the κ path.", {"size": 15.5, "space": 10}),
        ("The same modulus L powers the inversion-free certification of the "
         "hump region (next backup).", {"size": 14, "color": GREY}),
    ])

    # B5 certification method
    s = d.slide("Backup", "Certifying the hump region without inverting a Jacobian")
    d.body(s, Inches(0.85), Inches(1.95), Inches(11.65), Inches(1.6), [
        ("Threat: equilibrium cutoffs move with κ (GE channel B) and could "
         "overturn the fixed-cutoff hump (channel A).", {"size": 15,
         "space": 8}),
        ("Standard bound needs ‖dk/dκ‖ — the unknown path derivative.",
         {"size": 15, "space": 8}),
        ("Trick: F = I − T inherits the contraction modulus L of the "
         "best-response map, so by the Neumann series:", {"size": 15}),
    ])
    d.eq(s, "neumann", 16, cx=Inches(6.67), y=Inches(3.78))
    d.body(s, Inches(0.85), Inches(4.65), Inches(11.65), Inches(2.0), [
        ("— computable from **one solve plus one finite difference**.",
         {"size": 15, "space": 8}),
        ("Result: |B| bounded below channel A's amplitude on κ ∈ "
         "[0.35, 0.825] (inversion-free; [0.30, 0.85] with exact IFT). "
         "Margins: ball integral 2.5×10⁻⁵ vs 4.7×10⁻³.", {"size": 15,
         "space": 8}),
        ("At σ_ξ = 0.60 the same machinery certifies the **failure**: "
         "∫|B| = 0.021 vs A-amplitude 0.0069 → trough.", {"size": 14,
         "color": GREY}),
    ])

    # B6 tender game
    s = d.slide("Backup", "The tender game behind λ")
    d.body(s, Inches(0.85), Inches(1.92), Inches(11.65), Inches(1.1), [
        ("At the bargaining disagreement node: fringe raider (cost c_F ~ H, "
         "synergy S_F ~ G) posts a uniform tender for control share τ; "
         "atomistic float free-rides; charter allows dilution φ of "
         "non-tendering holders (Grossman–Hart).", {"size": 14.5}),
    ])
    d.body(s, Inches(0.85), Inches(3.00), Inches(7.0), Inches(0.5), [
        ("Free-rider floor — raid pays iff dilution on acquired shares "
         "covers entry cost:", {"size": 14.5}),
    ])
    d.eq(s, "freerider", 14, x=Inches(8.05), cy=Inches(3.28))
    d.body(s, Inches(0.85), Inches(3.85), Inches(7.0), Inches(0.5), [
        ("Pivotal bloc (1 − α < τ) can block: floor raids fail unless",
         {"size": 14.5}),
    ])
    d.eq(s, "block", 14, x=Inches(8.05), cy=Inches(4.01))
    d.body(s, Inches(0.85), Inches(4.70), Inches(7.0), Inches(0.5), [
        ("→ the ψ factor. The equilibrium share the activist's side keeps:",
         {"size": 14.5}),
    ])
    d.eq(s, "lam", 14, x=Inches(8.05), cy=Inches(4.86))
    d.hairline(s, Inches(0.85), Inches(5.42), Inches(11.65))
    d.body(s, Inches(0.85), Inches(5.60), Inches(11.65), Inches(1.2), [
        ("λ > 0 fails only if raids are certain, fully superseding, and "
         "unblockable — so the wedge m₁ > m₀ is a theorem on the rest of "
         "the parameter space.", {"size": 14.5}),
    ])

    # B7 sensitivity
    s = d.slide("Backup", "Sensitivity: the hump survives")
    d.picture(s, "sens1", MARGIN, Inches(2.0), w=Inches(6.1))
    d.picture(s, "sens2", Inches(6.85), Inches(2.0), w=Inches(6.1))
    d.caption(s, MARGIN, Inches(6.4), Inches(12),
              "Left: discount factor δ and bidder vol σ_ξ. Right: engagement cost C₀, "
              "premium wedge m₁−m₀, success probability ρ. Peak location κ† marked.")

    # B8 welfare
    s = d.slide("Backup", "Welfare vs minority gains")
    d.picture(s, "welfare", MARGIN, Inches(1.95), w=Inches(6.6))
    d.body(s, Inches(7.5), Inches(2.0), Inches(5.2), Inches(4.4), [
        ("Total surplus = synergies + improvements − costs; prices and "
         "premia are transfers.", {"size": 15, "space": 10}),
        ("The κ that maximizes minority extraction deters marginal "
         "efficient takeovers → κ* ≠ κ†.", {"size": 15, "space": 10}),
        ("Policy read: \"more liquidity is good for governance\" and \"good "
         "for welfare\" are different claims; the model separates them.",
         {"size": 14, "color": GREY}),
    ])

    # B9 disclosure benchmarks
    s = d.slide("Backup", "Disclosure benchmarks and rumors")
    d.picture(s, "rumor", MARGIN, Inches(1.95), w=Inches(6.6))
    d.body(s, Inches(7.5), Inches(2.0), Inches(5.2), Inches(4.4), [
        ("Bounds: full information (κ-invariant) and no disclosure "
         "(maximally κ-sensitive) bracket the baseline regime.",
         {"size": 15, "space": 10}),
        ("Noisy rumors (wolf-pack chatter) interpolate: as rumor precision "
         "rises, inference converges to full information and the liquidity "
         "effect flattens further.", {"size": 15, "space": 10}),
        ("Same comparative static as tightening the disclosure rule — "
         "information substitutes are interchangeable here.", {"size": 14,
         "color": GREY}),
    ])

    # B10 cutoffs vs kappa
    s = d.slide("Backup", "Equilibrium cutoffs across the liquidity dial")
    d.picture(s, "cutoffs_kappa", Inches(0.85), Inches(1.9), h=Inches(4.1))
    d.body(s, Inches(7.5), Inches(2.0), Inches(5.2), Inches(4.2), [
        ("Exit/hold boundary k₁ = k₀ throughout (hold collapsed); public "
         "threshold k_D tracks the quiet/public margin.", {"size": 15,
         "space": 10}),
        ("Ordering k₁ ≤ k₀ ≤ k_D preserved across all κ — the GE channel "
         "moves cutoffs smoothly, no regime jumps.", {"size": 15}),
    ])

    # B11 fact 2 methods
    s = d.slide("Backup", "Evidence: data & methods")
    d.body(s, Inches(0.85), Inches(2.0), Inches(11.65), Inches(4.6), [
        ("Universe: every original Schedule 13D on EDGAR 2022Q1–2025Q4 "
         "(amendments excluded; boundary-aware form matching — EDGAR renamed "
         "the form \"SCHEDULE 13D\" in Dec 2024).", {"size": 14.5,
         "space": 10}),
        ("CUSIPs parsed from cover pages with check-digit validation "
         "(~80–90% coverage); after-16:00 ET acceptances shifted to the next "
         "trading day.", {"size": 14.5, "space": 10}),
        ("CARs: market model, estimation 220 trading days ending 30 before "
         "the event (min 100 obs), windows [−1,+1] and [−10,+1].",
         {"size": 14.5, "space": 10}),
        ("Covariates: Amihud illiquidity [−250,−30], run-up [−60,−11], "
         "ln market cap at −30; US common stock only; two-way clustered SEs "
         "(firm × month); year-quarter FE absorb Post in the saturated spec.",
         {"size": 14.5}),
    ])

    d.prs.save(OUT_PPTX)
    print(f"wrote {OUT_PPTX}: {d.n} slides")


if __name__ == "__main__":
    build_assets()
    build()
